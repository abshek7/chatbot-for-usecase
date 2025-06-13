from flask import Flask, request, jsonify,send_from_directory,render_template
from flask_cors import CORS
from PyPDF2 import PdfReader
from docx import Document
import pandas as pd
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
from pptx import Presentation
from PIL import Image
import io

app = Flask(__name__, static_folder='static', template_folder='templates')

CORS(app, resources={r"/*": {"origins": ["http://localhost:8000", "http://127.0.0.1:8000"]}})
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# Extract text from PDF
def extract_text_from_pdf(file):
    pdf_reader = PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text
    return text

# Extract text from DOCX
def extract_text_from_docx(file):
    doc = Document(file)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

# Extract text from TXT
def extract_text_from_txt(file):
    return file.read().decode("utf-8")

# Extract text from CSV or XLSX
def extract_text_from_spreadsheet(file, file_type):
    if file_type == 'csv':
        df = pd.read_csv(file)
    else:
        df = pd.read_excel(file)
    return df.to_string(index=False)

# Extract text from PowerPoint
def extract_text_from_ppt(file):
    prs = Presentation(file)
    text = ""
    for slide_num, slide in enumerate(prs.slides, 1):
        text += f"\n--- Slide {slide_num} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
    return text

# Extract text from images using Gemini Vision
def extract_text_from_image(file):
    try:
        # Read image file
        image_bytes = file.read()
        file.seek(0)  # Reset file pointer for potential reuse
        
        # Convert to PIL Image
        image = Image.open(io.BytesIO(image_bytes))
        
        # Use Gemini model for image analysis
        model = genai.GenerativeModel('gemini-2.5-flash-preview-05-20')
        
        prompt = """
        Please analyze this image and extract all visible text content. 
        Also describe any important visual elements, charts, diagrams, or information 
        that might be relevant for question-answering purposes.
        Provide a comprehensive description that includes:
        1. All readable text
        2. Description of visual elements
        3. Any data or information presented in charts/graphs
        4. Overall context and meaning of the image
        """
        
        response = model.generate_content([prompt, image])
        return f"Image Analysis:\n{response.text}"
        
    except Exception as e:
        return f"Error processing image: {str(e)}"

# General file text extractor
def get_file_text(file):
    file_type = file.filename.split('.')[-1].lower()
    if file_type == 'pdf':
        return extract_text_from_pdf(file)
    elif file_type == 'docx':
        return extract_text_from_docx(file)
    elif file_type == 'txt':
        return extract_text_from_txt(file)
    elif file_type in ['csv', 'xlsx']:
        return extract_text_from_spreadsheet(file, file_type)
    elif file_type in ['ppt', 'pptx']:
        return extract_text_from_ppt(file)
    elif file_type in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp']:
        return extract_text_from_image(file)
    else:
        return ""  # Unsupported file types will return empty string

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    return text_splitter.split_text(text)

def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details, if the answer is not in
    provided context just say, "answer is not available in the context", don't provide the wrong answer\n\n
    Context:\n{context}\n
    Question:\n{question}\n
    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-2.5-flash-preview-05-20", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    return load_qa_chain(model, chain_type="stuff", prompt=prompt)

@app.route('/')
def index():
    return render_template('index.html')  # Not send_from_directory


@app.route('/health')
def health_check():
    return jsonify({"status": "healthy"})

@app.route('/upload-files', methods=['POST'])
def upload_files():
    if 'files' not in request.files:
        return jsonify({"success": False, "message": "No files provided"}), 400
    
    files = request.files.getlist('files')
    if not files:
        return jsonify({"success": False, "message": "No files selected"}), 400

    processed_files = 0
    all_text = ""

    for file in files:
        try:
            file_text = get_file_text(file)
            if file_text.strip():
                all_text += f"\n--- Content from {file.filename} ---\n{file_text}\n"
                processed_files += 1
        except Exception as e:
            return jsonify({"success": False, "message": f"Error processing {file.filename}: {str(e)}"}), 500

    if all_text.strip():
        text_chunks = get_text_chunks(all_text)
        get_vector_store(text_chunks)
        return jsonify({
            "success": True,
            "message": "Files processed successfully",
            "processed_files": processed_files
        })
    else:
        return jsonify({"success": False, "message": "No valid text could be extracted from any uploaded files"}), 400

@app.route('/ask-question', methods=['POST'])
def ask_question():
    data = request.get_json()
    if not data or 'question' not in data:
        return jsonify({"success": False, "message": "No question provided"}), 400

    question = data['question']
    
    try:
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
        new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
        docs = new_db.similarity_search(question)

        chain = get_conversational_chain()
        response = chain({"input_documents": docs, "question": question}, return_only_outputs=True)
        answer = response["output_text"].strip()

        # If the answer is not available in context, fallback to general chat model
        if answer.lower().startswith("answer is not available in the context") or not answer:
            general_model = genai.GenerativeModel('gemini-2.5-flash-preview-05-20')
            general_response = general_model.generate_content(question)
            return jsonify({
                "success": True,
                "answer": general_response.text.strip()
            })
        else:
            return jsonify({
                "success": True,
                "answer": answer
            })

    except Exception as e:
        return jsonify({"success": False, "message": f"Error processing question: {str(e)}"}), 500

#@app.route('/favicon.ico')
#def favicon():
#   return send_from_directory('.', 'favicon.ico', mimetype='image/vnd.microsoft.icon')

if __name__ == "__main__":
    app.run(debug=True, port=8000)
